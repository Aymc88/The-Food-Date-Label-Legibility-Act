#!/usr/bin/env python3
"""Rebuild the Food Date Label Legibility Act deck (17 slides, current script).

Slide order follows the July 2026 script:
  chips → lemonade → orange juice → creamy → natural → crunchy #1 → crunchy #2
  → bottom line → human cost → gap (AB 660) → odd one out → the fix → safe harbor
  → proven path → global precedent → closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image
import io, os, tempfile

DIR = "/Users/amandachen/Documents/Antigravity/Food Date Lable Legibility Act"
os.chdir(DIR)

# ── Colors ──
CA_GOLD = RGBColor(0xB8, 0x86, 0x0B)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_RED = RGBColor(0xC4, 0x1E, 0x3A)
SLIDE_BG = RGBColor(0xF8, 0xF6, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x80, 0x80, 0x90)
CARD_BG = RGBColor(0xF0, 0xEE, 0xE8)
GOLD_BORDER = RGBColor(0xB8, 0x86, 0x0B)
GREEN = RGBColor(0x2D, 0x7D, 0x46)
IMG_BG = RGBColor(0xF0, 0xED, 0xE6)

FONT_SERIF = "Georgia"
FONT_SANS = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank


def _to_rgb(img_path):
    """Ensure image is RGB (not RGBA) for PowerPoint compatibility."""
    im = Image.open(img_path)
    if im.mode == "RGBA":
        bg = Image.new("RGB", im.size, (240, 237, 230))
        bg.paste(im, mask=im.split()[3])
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        bg.save(tmp, format="PNG")
        tmp.close()
        return tmp.name
    return img_path


def add_bg(slide, color=SLIDE_BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def tb(slide, l, t, w, h, text, sz=18, color=TEXT_DARK, bold=False, italic=False,
       font=FONT_SANS, align=PP_ALIGN.CENTER):
    """Add a centered textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return box


def tb_multi(slide, l, t, w, h, lines, base_sz=18, base_col=TEXT_DARK):
    """Multi-line textbox. lines = [(text, {props}), ...]"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, props) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(props.get("sz", base_sz))
        p.font.color.rgb = props.get("color", base_col)
        p.font.bold = props.get("bold", False)
        p.font.italic = props.get("italic", False)
        p.font.name = props.get("font", FONT_SANS)
        p.alignment = props.get("align", PP_ALIGN.CENTER)
    return box


def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    """Rounded rectangle."""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    # Round corners
    spPr = s._element.spPr
    pG = spPr.find(qn('a:prstGeom'))
    if pG is None:
        pG = etree.SubElement(spPr, qn('a:prstGeom'))
    pG.set('prst', 'roundRect')
    return s


def img_slide(tag, heading, img, copy_img=None):
    """Full-bleed image slide with overlay."""
    sl = prs.slides.add_slide(blank)
    add_bg(sl, IMG_BG)

    # Background image — full size
    img_src = _to_rgb(img) if img.endswith('.png') else img
    sl.shapes.add_picture(img_src, 0, 0, SLIDE_W, SLIDE_H)

    if copy_img and os.path.exists(copy_img):
        copy_src = _to_rgb(copy_img)
        sl.shapes.add_picture(copy_src, 0, 0, SLIDE_W, SLIDE_H)

    # Overlay bar at bottom
    rect(sl, 0, 5.6, 13.333, 1.9, fill=RGBColor(0xF8, 0xF6, 0xF0))

    # Gold tag
    tag_s = rect(sl, 0.6, 5.8, 4.0, 0.4, fill=CA_GOLD, border=CA_GOLD)
    ttf = tag_s.text_frame
    ttf.paragraphs[0].text = tag
    ttf.paragraphs[0].font.size = Pt(12)
    ttf.paragraphs[0].font.color.rgb = WHITE
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.name = FONT_SANS
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Heading
    tb(sl, 0.6, 6.3, 12.0, 0.8, heading, sz=22, bold=True, color=TEXT_DARK,
       font=FONT_SERIF, align=PP_ALIGN.LEFT)
    return sl


def text_slide(badge, heading, build_fn=None):
    """Standard text slide."""
    sl = prs.slides.add_slide(blank)
    add_bg(sl, SLIDE_BG)

    if badge:
        b = rect(sl, 5.5, 0.4, 2.3, 0.4, border=CA_GOLD, bw=Pt(1.5))
        tf = b.text_frame
        tf.paragraphs[0].text = badge
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = CA_GOLD
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_SANS
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if heading:
        tb(sl, 1.0, 1.2, 11.333, 1.0, heading, sz=30, bold=True,
           color=TEXT_DARK, font=FONT_SERIF)

    if build_fn:
        build_fn(sl)
    return sl


# ── Build ──
print("Building slides...")

# 1 — Title
sl0 = prs.slides.add_slide(blank)
add_bg(sl0)
rect(sl0, 4.2, 1.8, 4.9, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl0, 4.2, 1.8, 4.9, 0.45, "PROPOSED LEGISLATION", sz=11, color=CA_GOLD, bold=True)
tb(sl0, 1.0, 2.8, 11.333, 2.0,
   "The Food Date Label\nLegibility Act", sz=48, bold=True, color=TEXT_DARK, font=FONT_SERIF)
tb(sl0, 1.0, 5.0, 11.333, 0.6, "By Amanda Chen", sz=18, color=CA_GOLD, bold=True)
print("  1/17 Title")

# 2 — Chips (overlapping)
img_slide("THE OVERLAPPING PROBLEM",
          "Can you read me the date on these chips?",
          "page_7.png")
print("  2/17 Chips")

# 3 — Lemonade (low contrast)
img_slide("THE LOW CONTRAST PROBLEM",
          "Clear embossing on clear plastic — the date is almost invisible.",
          "page_8.png", "page_8_copy.png")
print("  3/17 Lemonade")

# 4 — Orange juice (the exception)
img_slide("HIGH CONTRAST · THE EXCEPTION",
          "Easy to read. …But I don't even like orange juice.",
          "page_9.png", "page_9_copy.png")
print("  4/17 Orange Juice")

# 5–8 — Peanut butter sequence (same brand, same shelf, same store)
img_slide("SAME BRAND · SAME SHELF · SAME STORE",
          "Creamy peanut butter — I can't read the date.",
          "page_3.png", "page_3_copy.png")
print("  5/17 Creamy PB")

img_slide("SAME BRAND · SAME SHELF · SAME STORE",
          "Natural creamy peanut butter — I can't read the date.",
          "page_4.png", "page_4_copy.png")
print("  6/17 Natural PB")

img_slide("SAME BRAND · SAME SHELF · SAME STORE",
          "Crunchy #1 — I can't read the date.",
          "page_5.png", "page_5_copy.png")
print("  7/17 Crunchy #1")

img_slide("THE DISCREPANCY PROBLEM",
          "Crunchy #2 — clear as day. One jar legible; the other faint carvings on plastic.",
          "page_6.png", "page_6_copy.png")
print("  8/17 Crunchy #2")

# 9 — Bottom line
sl9 = prs.slides.add_slide(blank)
add_bg(sl9)
rect(sl9, 5.0, 0.5, 3.3, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl9, 5.0, 0.5, 3.3, 0.45, "THE BOTTOM LINE", sz=11, color=CA_GOLD, bold=True)
tb(sl9, 1.0, 1.3, 11.333, 0.8, "From just one routine grocery run…",
   sz=30, bold=True, color=TEXT_DARK, font=FONT_SERIF)
for i, (num, label) in enumerate([
    ("4/7", "Failure rate — unreadable dates"),
    ("2/3", "Forced compromise — not my first choice"),
    ("8/25", "Camera test — missed in a standard pantry"),
]):
    x = 1.5 + i * 3.7
    rect(sl9, x, 3.0, 3.2, 2.2, fill=CARD_BG, border=RGBColor(0xD0, 0xD0, 0xD0))
    tb(sl9, x, 3.3, 3.2, 0.9, num, sz=48, bold=True,
       color=ACCENT_RED if i != 1 else CA_GOLD)
    tb(sl9, x + 0.2, 4.3, 2.8, 0.8, label, sz=13, color=TEXT_DARK)
tb(sl9, 1.0, 5.6, 11.333, 0.7,
   "Good eyesight. No crowd. No time pressure. Perfect lighting. — It was still hard to read.",
   sz=15, italic=True, color=MUTED)
print("  9/17 Bottom Line")

# 10 — Human cost
sl10 = prs.slides.add_slide(blank)
add_bg(sl10)
rect(sl10, 4.9, 0.5, 3.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl10, 4.9, 0.5, 3.5, 0.45, "THE HUMAN COST", sz=11, color=CA_GOLD, bold=True)
tb(sl10, 1.0, 1.2, 11.333, 1.0,
   "When shoppers can't read a date label,\nthey do one of two things:",
   sz=28, bold=True, color=TEXT_DARK, font=FONT_SERIF)
for i, opt in enumerate(["Eat food they shouldn't", "Throw out food that was perfectly fine"]):
    x = 2.0 + i * 5.0
    rect(sl10, x, 3.0, 4.2, 1.5, fill=RGBColor(0xFD, 0xF6, 0xE0), border=GOLD_BORDER, bw=Pt(1.5))
    tb(sl10, x, 3.0, 4.2, 1.5, opt, sz=18, bold=True, color=TEXT_DARK)
tb_multi(sl10, 2.0, 5.0, 9.333, 1.2, [
    ("3 billion pounds of food · $7 billion — every single year.", {"sz": 20, "bold": True}),
    ("We fixed the words — now let's make them readable.", {"sz": 18, "bold": True, "color": CA_GOLD}),
])
print("  10/17 Human Cost")

# 11 — The Gap (AB 660)
sl11 = prs.slides.add_slide(blank)
add_bg(sl11)
rect(sl11, 5.9, 0.5, 1.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl11, 5.9, 0.5, 1.5, 0.45, "THE GAP", sz=11, color=CA_GOLD, bold=True)
tb_multi(sl11, 1.0, 1.2, 11.333, 1.5, [
    ("AB 660 fixed what to say.", {"sz": 30, "bold": True}),
    ("It never said make it readable.", {"sz": 24, "bold": True, "color": ACCENT_RED}),
])
rows = [("Standardized words — “Best if Used by,” “Use by”", "✓ AB 660 (2024)", GREEN),
        ("Confusing “Sell by” dates", "✓ Banned", GREEN),
        ("Legibility — can you read it?", "✗ No requirement", ACCENT_RED)]
for i, (lt, rt, rc) in enumerate(rows):
    y = 3.4 + i * 0.75
    bg_c = RGBColor(0xF5, 0xF3, 0xEE) if i < 2 else RGBColor(0xFD, 0xE8, 0xE8)
    bd_c = RGBColor(0xDD, 0xDD, 0xDD) if i < 2 else RGBColor(0xF0, 0xC0, 0xC0)
    rect(sl11, 2.5, y, 8.3, 0.6, fill=bg_c, border=bd_c)
    tb(sl11, 2.8, y, 4.4, 0.6, lt, sz=14, bold=(i == 2), align=PP_ALIGN.LEFT)
    tb(sl11, 7.2, y, 3.4, 0.6, rt, sz=14, bold=True, color=rc, align=PP_ALIGN.RIGHT)
tb(sl11, 1.0, 5.7, 11.333, 0.9,
   "A company can print the exact right phrase — in ink the same color as the background,\nsqueezed under the barcode, or stamped over its own logo — and it's fully legal.",
   sz=14, color=MUTED)
print("  11/17 The Gap")

# 12 — The Odd One Out
sl12 = prs.slides.add_slide(blank)
add_bg(sl12)
rect(sl12, 4.9, 0.5, 3.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl12, 4.9, 0.5, 3.5, 0.45, "THE ODD ONE OUT", sz=11, color=CA_GOLD, bold=True)
tb(sl12, 1.0, 1.2, 11.333, 0.8, "Everything else on the package must be legible.",
   sz=26, bold=True, color=TEXT_DARK, font=FONT_SERIF)
rows = [("Nutrition Facts panel", "✓ “Plainly legible” — federal law", GREEN),
        ("Net quantity statement", "✓ “Distinct contrast”", GREEN),
        ("Alcohol labels", "✓ Contrast rule", GREEN),
        ("The date", "✗ No legibility standard at all", ACCENT_RED)]
for i, (lt, rt, rc) in enumerate(rows):
    y = 2.7 + i * 0.8
    bg_c = RGBColor(0xF5, 0xF3, 0xEE) if i < 3 else RGBColor(0xFD, 0xE8, 0xE8)
    bd_c = RGBColor(0xDD, 0xDD, 0xDD) if i < 3 else RGBColor(0xF0, 0xC0, 0xC0)
    rect(sl12, 2.5, y, 8.3, 0.65, fill=bg_c, border=bd_c)
    tb(sl12, 2.8, y, 4.4, 0.65, lt, sz=15, bold=(i == 3), align=PP_ALIGN.LEFT)
    tb(sl12, 7.2, y, 3.4, 0.65, rt, sz=14, bold=True, color=rc, align=PP_ALIGN.RIGHT)
tb(sl12, 1.0, 6.0, 11.333, 0.8,
   "The calorie count has to be legible. The net weight has to be legible.\nThe date is the odd one out.",
   sz=16, bold=True, color=CA_GOLD)
print("  12/17 The Odd One Out")

# 13 — The Fix
sl13 = prs.slides.add_slide(blank)
add_bg(sl13)
rect(sl13, 5.8, 0.3, 1.7, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl13, 5.8, 0.3, 1.7, 0.45, "THE FIX", sz=11, color=CA_GOLD, bold=True)
tb(sl13, 1.0, 1.0, 11.333, 0.6, "My bill adds one section to existing law.", sz=26, bold=True, color=TEXT_DARK, font=FONT_SERIF)
tb(sl13, 1.0, 1.6, 11.333, 0.4, "Three simple requirements.", sz=18, bold=True, color=CA_GOLD)
for i, (num, title, desc) in enumerate([
    ("01", "Contrast", "4.5:1 contrast ratio — the same standard accessibility experts already use. Black on white passes automatically: a safe harbor, not a mandate."),
    ("02", "Clear Field", "No other text, graphic, or logo overlapping the date. The date can't be buried under marketing."),
    ("03", "Forward-Looking", "Only applies to products made after July 2029. Two years of runway. No company has to throw away packaging they've already printed."),
]):
    x = 0.8 + i * 4.1
    rect(sl13, x, 2.5, 3.6, 3.8, fill=CARD_BG, border=RGBColor(0xD0, 0xD0, 0xD0))
    tb(sl13, x, 2.7, 3.6, 0.5, num, sz=32, bold=True, color=CA_GOLD, align=PP_ALIGN.LEFT)
    tb(sl13, x + 0.2, 3.4, 3.2, 0.4, title, sz=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT)
    tb(sl13, x + 0.2, 4.0, 3.2, 2.0, desc, sz=12, color=MUTED, align=PP_ALIGN.LEFT)
tb(sl13, 1.0, 6.5, 11.333, 0.6,
   "No new agency. No new paperwork. It folds into existing enforcement — it costs the state essentially nothing.",
   sz=14, bold=True, color=CA_GOLD)
print("  13/17 The Fix")

# 14 — Safe Harbor
sl14 = prs.slides.add_slide(blank)
add_bg(sl14)
rect(sl14, 3.4, 0.4, 6.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl14, 3.4, 0.4, 6.5, 0.45, "SAFE HARBOR, NOT A MANDATE", sz=11, color=CA_GOLD, bold=True)
for i, (ql, qt, ans) in enumerate([
    ("Question #1", "Isn't a formatting rule too prescriptive?",
     "4.5:1 is a safe harbor, not a mandate. If a manufacturer prints black ink on a white box, they automatically comply — no calculations needed."),
    ("Question #2", "Isn't a contrast ratio a screen standard?",
     "I'm borrowing the number, not the screen. Contrast is just math about two colors — true on a phone or a peanut butter jar. The bill lets the Department of Food and Agriculture set the measurement method for physical packages."),
]):
    x = 0.8 + i * 6.5
    rect(sl14, x, 1.5, 5.7, 4.5, fill=RGBColor(0xFD, 0xF6, 0xE0), border=GOLD_BORDER, bw=Pt(1.5))
    tb(sl14, x + 0.3, 1.8, 3.0, 0.3, ql, sz=11, bold=True, color=CA_GOLD, align=PP_ALIGN.LEFT)
    tb(sl14, x + 0.3, 2.2, 5.1, 0.6, qt, sz=16, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT)
    tb(sl14, x + 0.3, 2.9, 5.1, 2.5, ans, sz=13, color=MUTED, align=PP_ALIGN.LEFT)
print("  14/17 Safe Harbor")

# 15 — A Proven Path
sl15 = prs.slides.add_slide(blank)
add_bg(sl15)
rect(sl15, 4.9, 0.5, 3.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl15, 4.9, 0.5, 3.5, 0.45, "A PROVEN PATH", sz=11, color=CA_GOLD, bold=True)
tb(sl15, 1.0, 1.3, 11.333, 0.8, "California didn't get here in one shot.",
   sz=28, bold=True, color=TEXT_DARK, font=FONT_SERIF)
for i, (year, label, highlight) in enumerate([
    ("2016", "A bill that failed", False),
    ("2017", "Voluntary version", False),
    ("2024", "AB 660 — standard words", False),
    ("2029", "Legible dates — the next step", True),
]):
    x = 1.2 + i * 2.9
    if highlight:
        rect(sl15, x, 2.7, 2.4, 2.2, fill=RGBColor(0xFD, 0xF6, 0xE0), border=GOLD_BORDER, bw=Pt(2))
        tb(sl15, x, 3.0, 2.4, 0.6, year, sz=30, bold=True, color=CA_GOLD)
        tb(sl15, x + 0.15, 3.7, 2.1, 1.0, label, sz=13, bold=True, color=TEXT_DARK)
    else:
        rect(sl15, x, 2.7, 2.4, 2.2, fill=CARD_BG, border=RGBColor(0xD0, 0xD0, 0xD0))
        tb(sl15, x, 3.0, 2.4, 0.6, year, sz=30, bold=True, color=ACCENT_RED)
        tb(sl15, x + 0.15, 3.7, 2.1, 1.0, label, sz=13, color=MUTED)
tb(sl15, 1.0, 5.4, 11.333, 0.8,
   "We fixed the words — now we make them legible.",
   sz=20, bold=True, color=CA_GOLD)
print("  15/17 Proven Path")

# 16 — Global Precedent
sl16 = prs.slides.add_slide(blank)
add_bg(sl16)
rect(sl16, 3.4, 0.5, 6.5, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl16, 3.4, 0.5, 6.5, 0.45, "NOT RADICAL — IT'S NORMAL EVERYWHERE ELSE", sz=11, color=CA_GOLD, bold=True)
tb_multi(sl16, 1.5, 1.9, 10.333, 1.8, [
    ("The European Union requires date information to be", {"sz": 20}),
    ('"clearly legible" and — almost word for word what I\'m proposing —', {"sz": 20}),
    ('"not hidden, obscured, or interrupted by any other matter."', {"sz": 20}),
])
tb(sl16, 1.5, 4.0, 10.333, 0.6, "Canada requires it too.", sz=20, bold=True, color=TEXT_DARK)
tb(sl16, 1.0, 5.0, 11.333, 0.9,
   "California led the country on standardizing the words.\nWe can lead on making them readable.",
   sz=20, bold=True, color=CA_GOLD)
print("  16/17 Global Precedent")

# 17 — Closing
sl17 = prs.slides.add_slide(blank)
add_bg(sl17)
rect(sl17, 4.2, 1.2, 4.9, 0.45, border=CA_GOLD, bw=Pt(1.5))
tb(sl17, 4.2, 1.2, 4.9, 0.45, "PROPOSED LEGISLATION", sz=11, color=CA_GOLD, bold=True)
tb(sl17, 1.0, 2.0, 11.333, 1.5,
   "The Food Date Label\nLegibility Act", sz=44, bold=True, color=TEXT_DARK, font=FONT_SERIF)
tb(sl17, 1.0, 4.0, 11.333, 1.0,
   "California already tells food companies what these dates must say.\nIt just never required them to be readable.",
   sz=16, color=MUTED)
tb(sl17, 1.0, 5.4, 11.333, 0.7,
   "We fixed the words — now let's make them readable.",
   sz=22, bold=True, color=CA_GOLD)
print("  17/17 Closing")

# ── Speaker notes ──
NOTES = [
    "Opening:\nGood Morning/Afternoon, Assemblymember Lee and other legislative representatives. In California, food companies are legally required to put 'Best By' dates on packages, yet millions of shoppers can't actually read them. In the next few minutes I will show you how this oversight forces families to throw away billions in good food, and how three simple standards—contrast, location, and a clear field—will fix it without costing taxpayers a single cent. This is The Food Date Label Legibility Act.",
    "Chips:\nAssemblymember Lee, please look at this bag of chips. Every item is required by law to carry a date label, but can you read me the date on those chips? [7 seconds] Hardly. The silver ink is stamped directly over the light brand logo. When marketing and safety dates compete for space, marketing wins every time. That's the overlapping problem.",
    "Lemonade:\nNext, I take your attention to the screen. This is my lemonade. Can anyone read me the date? [3 seconds] No, you cannot. The clear embossing on clear plastic makes the date almost invisible. That's the low-contrast problem.",
    "Orange juice:\nSo for safety, I switched to orange juice—high contrast, easy to read. …But I don't even like orange juice.",
    "Creamy PB:\nAnd here's my peanut butter—I actually ate this for breakfast today—same brand, same shelf, same store. Look at creamy peanut butter. I can't read the date.",
    "Natural PB:\nNatural creamy peanut butter. I can't read the date.",
    "Crunchy #1:\nCrunchy peanut butter number 1. I can't read the date.",
    "Crunchy #2:\nCrunchy peanut butter number 2. Clear as day. One jar is clear and legible; the other is faint carvings on plastic. That's the discrepancy problem. So for safety, I switched to Crunchy #2—high contrast, easy to read. …But I don't even like crunchy peanut butter.",
    "Bottom line:\nHere is the bottom line from just one routine grocery run. First, failure rate: more than half the items I picked up—4 out of 7—had unreadable dates. Second, forced compromise: of the 3 things I actually bought, 2 were not my first choice—not because of price, and not because of taste, but purely because the dates were invisible. Third, camera test: when I tested 25 everyday household food items under optimal camera lighting, my phone camera almost missed 8 of them—nearly a third of a standard pantry. Now—you all have good eyesight. You had no crowded line behind you, no time pressure, and perfect room lighting—and it was still hard to read. Picture my highly-nearsighted mom in a real California grocery aisle. She's under harsh fluorescent lights. Behind her, someone's rushing past with a cart, an auntie in a rush to work is ushering her along, and she's squinting through her glasses, trying in vain to decipher gray ink on a brown lid.",
    "Human cost:\nWhen shoppers can't read a date label, they do one of two things: they eat food they shouldn't, or—far more often—they throw out food that is perfectly fine. That confusion drives Americans to throw away 3 billion pounds of food—worth 7 billion dollars—every single year. California already tells food companies what these dates must say. It just never required them to be readable. We fixed the words—now let's make them readable. My bill closes that gap.",
    "The gap (AB 660):\nTwo years ago, California passed AB 660. Good law. It said: food companies have to use clear, standardized date words—'Best if Used by,' 'Use by'—and it banned the confusing 'sell by' dates. It went into effect this July first. But AB 660 only fixed the words. It said nothing about whether you can actually read them. So a company can print the exact right phrase—in ink the same color as the background, or squeezed under the barcode, or stamped over its own logo—and it's fully legal. The law tells them what to say. It never says make it readable. That's the gap.",
    "Odd one out:\nCalifornia and the federal government already require legibility for almost everything else on that same package. The Nutrition Facts panel? Federal law requires it in a type size and color contrast that's—and I quote—'plainly legible.' The net-quantity statement? Required to appear in 'distinct contrast' to everything around it. Even alcohol labels have a contrast rule. So on a single package, the calorie count has to be legible. The net weight has to be legible. But the date—the one piece of information that tells you if the food is still safe—is the one thing with no legibility standard at all. The date is the odd one out. My bill doesn't invent a new kind of rule. It closes an obvious gap where the date got left behind.",
    "The fix:\nSo here's what my bill actually does. It adds one section to the existing date-label law. Three simple requirements. One, contrast: the date has to stand out from its background enough to be read. To make that objective I use the same standard accessibility experts already use: a 4.5-to-1 contrast ratio. If you print black on white, you automatically pass—that's a safe harbor, not a mandate. On some of the products I showed you today, the printing was so faint it was nearly the same color as the background, almost invisible—approaching a 1-to-1 ratio. This is far below the 4.5-to-1 minimum of WCAG 2.1, a guideline that Section 508 of the Rehabilitation Act and ADA Title II require. That's not just a readability issue—that's an invisibility issue. Two, a clear field: no other text, graphic, or logo overlapping the date. Three, it only applies going forward, to products made after July 2029—two years of runway. No new agency. No new paperwork. It folds into the enforcement that already exists. It costs the state essentially nothing.",
    "Safe harbor:\nNow—I know the pushback. Food manufacturers will say a formatting rule is too prescriptive. Again: 4.5-to-1 is a safe harbor, not a mandate. If a manufacturer prints black ink on a white box, they automatically comply—no calculations needed. They might claim a contrast ratio built for computer screens doesn't belong on a package. Fair concern. I'm borrowing the number, not the screen. A contrast ratio is just math about two colors—it's true whether it's on a phone or a peanut butter jar. My bill lets the Department of Food and Agriculture set the exact measurement method for physical packages.",
    "Proven path:\nThis is also a proven path. California didn't standardize date labels in one shot—it took three tries over nearly a decade. A bill in 2016 that failed. A voluntary version in 2017. And finally AB 660 in 2024. Each one built on the last. My bill is simply the next step: we fixed the words; now we make them legible.",
    "Global precedent:\nAnd this isn't radical—it's normal everywhere else. The European Union requires date information to be 'clearly legible' and—almost word for word what I'm proposing—'not hidden, obscured, or interrupted by any other matter.' Canada requires it too. California led the country on standardizing the words. We can lead on making them readable.",
    "Closing:\nCalifornia already tells food companies what these dates must say. It just never required them to be readable. This is The Food Date Label Legibility Act. We fixed the words—now let's make them readable.",
]

for i, sl in enumerate(prs.slides):
    try:
        ns = sl.notes_slide
        ns.notes_text_frame.text = NOTES[i] if i < len(NOTES) else ""
    except:
        pass

# ── Save ──
out = os.path.join(DIR, "The Food Date Label Legibility Act - PPT.pptx")
prs.save(out)
print(f"\nSaved: {out}")
print(f"Size: {os.path.getsize(out) / 1024 / 1024:.1f} MB")
